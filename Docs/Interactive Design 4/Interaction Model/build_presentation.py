"""
Build: The Interaction Atom — presentation for ID4 students.

Style matches existing course presentations:
  - Dark background, white text
  - Gold accent for labels, blue for secondary
  - One idea per slide, large type
  - 10" x 5.625" (16:9)
"""

import math
import os
import tempfile
import numpy as np
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ═══════════════════════════════════════════════════════════════
#  STYLE CONSTANTS
# ═══════════════════════════════════════════════════════════════

# Slide size
SLIDE_W = Inches(10)
SLIDE_H = Inches(5.625)

# Colors
BG_COLOR     = RGBColor(0x1A, 0x1A, 0x2E)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY   = RGBColor(0xB0, 0xB0, 0xB0)
MID_GRAY     = RGBColor(0x80, 0x80, 0x80)
DARK_GRAY    = RGBColor(0x55, 0x55, 0x55)
GOLD         = RGBColor(0xD4, 0xA8, 0x53)
BLUE         = RGBColor(0x6B, 0x9E, 0xBF)
RED          = RGBColor(0xE7, 0x4C, 0x3C)
GREEN        = RGBColor(0x27, 0xAE, 0x60)
AMBER        = RGBColor(0xF3, 0x9C, 0x12)

# Mode colors
MODE_COLORS = {
    "Tracking":          BLUE,
    "Switch":            AMBER,
    "One-Shot":          RED,
    "Threshold Trigger": GREEN,
}

MODE_BG = {
    "Tracking":          RGBColor(0x1E, 0x3A, 0x5F),
    "Switch":            RGBColor(0x4A, 0x3A, 0x1A),
    "One-Shot":          RGBColor(0x5A, 0x1A, 0x1A),
    "Threshold Trigger": RGBColor(0x1A, 0x4A, 0x2A),
}


# ═══════════════════════════════════════════════════════════════
#  HELPERS
# ═══════════════════════════════════════════════════════════════

def set_slide_bg(slide, color=BG_COLOR):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text,
                font_size=14, color=WHITE, bold=False, italic=False,
                alignment=PP_ALIGN.LEFT, font_name="Calibri",
                line_spacing=None, anchor=MSO_ANCHOR.TOP):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None

    # Handle multi-line with format markers
    lines = text.split("\n")
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = alignment

        if line_spacing:
            p.line_spacing = line_spacing

        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.bold = bold
        run.font.italic = italic
        run.font.name = font_name

    return txBox


def add_rich_text(slide, left, top, width, height, segments,
                  alignment=PP_ALIGN.LEFT, line_spacing=None):
    """Add textbox with mixed formatting.
    segments = list of (text, font_size, color, bold, italic) tuples.
    Use '\\n' in text for new paragraphs.
    """
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.alignment = alignment
    if line_spacing:
        p.line_spacing = line_spacing

    for text, fsize, color, bold, italic in segments:
        parts = text.split("\n")
        for j, part in enumerate(parts):
            if j > 0:
                p = tf.add_paragraph()
                p.alignment = alignment
                if line_spacing:
                    p.line_spacing = line_spacing
            run = p.add_run()
            run.text = part
            run.font.size = Pt(fsize)
            run.font.color.rgb = color
            run.font.bold = bold
            run.font.italic = italic
            run.font.name = "Calibri"

    return txBox


def add_rect(slide, left, top, width, height, fill_color, border_color=None, border_width=Pt(1)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_arrow_line(slide, x1, y1, x2, y2, color=LIGHT_GRAY, width=Pt(2)):
    connector = slide.shapes.add_connector(
        1, x1, y1, x2, y2  # 1 = straight connector
    )
    connector.line.color.rgb = color
    connector.line.width = width
    return connector


def section_label(slide, text):
    """Small gold label at top-left."""
    add_textbox(slide, Inches(0.6), Inches(0.3), Inches(4), Inches(0.4),
                text, font_size=13, color=GOLD, bold=True)


def slide_title(slide, text, font_size=32):
    """Large white title."""
    add_textbox(slide, Inches(0.6), Inches(0.8), Inches(8.8), Inches(1.2),
                text, font_size=font_size, color=WHITE, bold=True)


def slide_body(slide, text, top=Inches(2.2), font_size=14):
    """Body text in light gray."""
    add_textbox(slide, Inches(0.6), top, Inches(8.8), Inches(3),
                text, font_size=font_size, color=LIGHT_GRAY,
                line_spacing=Pt(22))


# ═══════════════════════════════════════════════════════════════
#  CURVE PLOTS (for embedding)
# ═══════════════════════════════════════════════════════════════

x_vals = np.linspace(0, 1, 200)

CURVES = {
    "Linear":      lambda x: x,
    "Ease-in":     lambda x: x**2,
    "Ease-out":    lambda x: 1 - (1-x)**2,
    "Ease-in-out": lambda x: np.where(x < 0.5, 2*x**2, 1-(-2*x+2)**2/2),
    "Inverse":     lambda x: 1 - x,
    "S-curve":     lambda x: 3*x**2 - 2*x**3,
}

RESPONSE_CURVES = {
    "Burst":     lambda x: np.where(x <= 0, 0, np.minimum(1.5, (x/0.1)*np.exp(1-x/0.1))),
    "Swell":     lambda x: np.sin(np.pi * x),
    "Snap+tail": lambda x: np.exp(-4*x),
    "Overshoot": lambda x: np.where(x <= 0, 0, np.minimum(1.5, 1.4*(x/0.15)*np.exp(1-x/0.15))),
    "Spike":     lambda x: np.where(x < 0.05, x/0.05, np.maximum(0, (1-x)/0.95)),
}


def make_curve_strip(curves_dict, color, title, y_max=1.15):
    """Generate a horizontal strip of small curve plots, return temp PNG path."""
    n = len(curves_dict)
    fig, axes = plt.subplots(1, n, figsize=(n * 1.6, 1.4))
    fig.patch.set_facecolor("#1A1A2E")

    if n == 1:
        axes = [axes]

    for ax, (name, func) in zip(axes, curves_dict.items()):
        y = func(x_vals)
        ax.plot(x_vals, y, color=color, linewidth=2)
        ax.axhline(y=1, color="#444444", linewidth=0.5, linestyle="--")
        ax.set_xlim(0, 1)
        ax.set_ylim(-0.05, y_max)
        ax.set_title(name, fontsize=8, color="white", pad=3, fontfamily="sans-serif")
        ax.set_xticks([])
        ax.set_yticks([])
        ax.set_facecolor("#1A1A2E")
        for spine in ax.spines.values():
            spine.set_color("#444444")
            spine.set_linewidth(0.5)

    fig.tight_layout(pad=0.3)
    tmp = tempfile.NamedTemporaryFile(suffix=".png", delete=False)
    fig.savefig(tmp.name, dpi=200, bbox_inches="tight", facecolor="#1A1A2E")
    plt.close(fig)
    return tmp.name


def make_comparison_plot():
    """3 curves side by side showing same input/output with different feel."""
    fig, axes = plt.subplots(1, 3, figsize=(6, 1.8))
    fig.patch.set_facecolor("#1A1A2E")

    configs = [
        ("Linear", CURVES["Linear"], "#6B9EBF", "Attentive, crisp"),
        ("Ease-out", CURVES["Ease-out"], "#D4A853", "Responsive, sticky"),
        ("Inverse", CURVES["Inverse"], "#E74C3C", "Reluctant, evasive"),
    ]

    for ax, (name, func, color, feel) in zip(axes, configs):
        y = func(x_vals)
        ax.plot(x_vals, y, color=color, linewidth=2.5)
        ax.axhline(y=1, color="#444444", linewidth=0.5, linestyle="--")
        ax.set_xlim(0, 1)
        ax.set_ylim(-0.05, 1.15)
        ax.set_title(name, fontsize=10, color="white", pad=4, fontweight="bold",
                     fontfamily="sans-serif")
        ax.set_xlabel(feel, fontsize=7, color="#888888", labelpad=2, fontfamily="sans-serif")
        ax.set_xticks([0, 1])
        ax.set_xticklabels(["near", "far"], fontsize=6, color="#666666")
        ax.set_yticks([0, 1])
        ax.set_yticklabels(["dim", "bright"], fontsize=6, color="#666666")
        ax.set_facecolor("#1A1A2E")
        for spine in ax.spines.values():
            spine.set_color("#444444")
            spine.set_linewidth(0.5)

    fig.tight_layout(pad=0.5)
    tmp = tempfile.NamedTemporaryFile(suffix=".png", delete=False)
    fig.savefig(tmp.name, dpi=200, bbox_inches="tight", facecolor="#1A1A2E")
    plt.close(fig)
    return tmp.name


# ═══════════════════════════════════════════════════════════════
#  BUILD PRESENTATION
# ═══════════════════════════════════════════════════════════════

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# Use blank layout
blank_layout = prs.slide_layouts[6]  # Blank


# ── SLIDE 1: Title ──────────────────────────────────────────

s = prs.slides.add_slide(blank_layout)
set_slide_bg(s)
add_textbox(s, Inches(0.6), Inches(1.2), Inches(8.8), Inches(1.5),
            "THE INTERACTION", font_size=48, bold=True)
add_textbox(s, Inches(0.6), Inches(2.2), Inches(8.8), Inches(1.5),
            "ATOM", font_size=48, bold=True)
add_textbox(s, Inches(0.6), Inches(3.5), Inches(8.8), Inches(0.5),
            "Input  \u00b7  Transformation  \u00b7  Output", font_size=16, color=LIGHT_GRAY)
add_textbox(s, Inches(0.6), Inches(4.5), Inches(8.8), Inches(0.5),
            "Interaction Design 4  \u00b7  NABA 2026", font_size=11, color=MID_GRAY)


# ── SLIDE 2: The Question ───────────────────────────────────

s = prs.slides.add_slide(blank_layout)
set_slide_bg(s)
add_textbox(s, Inches(0.8), Inches(1.3), Inches(8.4), Inches(1.5),
            "When you walk toward a light", font_size=38, bold=True)
add_textbox(s, Inches(0.8), Inches(2.6), Inches(8.4), Inches(1.5),
            "and it gets brighter...", font_size=38, bold=True)
add_textbox(s, Inches(0.8), Inches(4.0), Inches(8.4), Inches(0.8),
            "what\u2019s actually happening?", font_size=16, color=LIGHT_GRAY)


# ── SLIDE 3: The Atom ───────────────────────────────────────

s = prs.slides.add_slide(blank_layout)
set_slide_bg(s)
section_label(s, "THE CORE MODEL")
slide_title(s, "Three parts. Every time.", font_size=28)

# Diagram: INPUT → TRANSFORMATION → OUTPUT
box_y = Inches(2.5)
box_h = Inches(1.4)

# INPUT box
add_rect(s, Inches(0.8), box_y, Inches(2.0), box_h, RGBColor(0x2C, 0x3E, 0x50), BLUE)
add_textbox(s, Inches(0.8), Inches(2.6), Inches(2.0), Inches(0.5),
            "INPUT", font_size=16, bold=True, alignment=PP_ALIGN.CENTER)
add_textbox(s, Inches(0.8), Inches(3.1), Inches(2.0), Inches(0.6),
            "what the space senses", font_size=10, color=LIGHT_GRAY,
            alignment=PP_ALIGN.CENTER)

# Arrow 1
add_textbox(s, Inches(2.9), Inches(2.85), Inches(0.5), Inches(0.5),
            "\u2192", font_size=24, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

# TRANSFORMATION box (larger, emphasized)
add_rect(s, Inches(3.5), box_y, Inches(3.2), box_h, RGBColor(0x2A, 0x2A, 0x45), GOLD, Pt(2))
add_textbox(s, Inches(3.5), Inches(2.55), Inches(3.2), Inches(0.5),
            "TRANSFORMATION", font_size=18, bold=True, color=GOLD,
            alignment=PP_ALIGN.CENTER)
add_textbox(s, Inches(3.5), Inches(3.05), Inches(3.2), Inches(0.7),
            "how input becomes output\nthis is where design lives",
            font_size=10, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Arrow 2
add_textbox(s, Inches(6.8), Inches(2.85), Inches(0.5), Inches(0.5),
            "\u2192", font_size=24, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

# OUTPUT box
add_rect(s, Inches(7.3), box_y, Inches(2.0), box_h, RGBColor(0x2C, 0x3E, 0x50), BLUE)
add_textbox(s, Inches(7.3), Inches(2.6), Inches(2.0), Inches(0.5),
            "OUTPUT", font_size=16, bold=True, alignment=PP_ALIGN.CENTER)
add_textbox(s, Inches(7.3), Inches(3.1), Inches(2.0), Inches(0.6),
            "what the space changes", font_size=10, color=LIGHT_GRAY,
            alignment=PP_ALIGN.CENTER)

add_textbox(s, Inches(0.8), Inches(4.5), Inches(8.4), Inches(0.5),
            "Every interaction in a responsive space is one atom.",
            font_size=13, color=MID_GRAY, italic=True, alignment=PP_ALIGN.CENTER)


# ── SLIDE 4: Input ──────────────────────────────────────────

s = prs.slides.add_slide(blank_layout)
set_slide_bg(s)
section_label(s, "INPUT")
slide_title(s, "What the space can sense about you")

inputs_text = (
    "SPATIAL\n"
    "Proximity \u2014 how far you are\n"
    "Gaze \u2014 where you\u2019re looking, for how long\n"
    "Stillness \u2014 how long since you moved\n"
    "Velocity \u2014 how fast you\u2019re moving\n"
    "\n"
    "DIRECT\n"
    "Contact \u2014 stepping into a zone\n"
    "Grab \u2014 picking up an object\n"
    "\n"
    "INTERFACE\n"
    "Button \u2014 press / release\n"
    "Slider \u2014 a value you control\n"
    "\n"
    "AUTONOMOUS\n"
    "Time \u2014 clock, breathing cycle\n"
    "Random \u2014 unpredictable variation"
)

# Split into columns
col1 = ("SPATIAL\n"
        "Proximity \u2014 how far you are\n"
        "Gaze \u2014 where you look, for how long\n"
        "Stillness \u2014 seconds since last movement\n"
        "Velocity \u2014 how fast you move")

col2 = ("DIRECT\n"
        "Contact \u2014 stepping into a zone\n"
        "Grab \u2014 picking up an object\n"
        "\n"
        "INTERFACE\n"
        "Button \u2014 press / release\n"
        "Slider \u2014 a continuous value")

col3 = ("AUTONOMOUS\n"
        "Time \u2014 clock, oscillation\n"
        "Random \u2014 unpredictable variation")

for i, (col_text, left) in enumerate([(col1, 0.6), (col2, 3.8), (col3, 7.0)]):
    lines = col_text.split("\n")
    txBox = s.shapes.add_textbox(Inches(left), Inches(2.0), Inches(3.0), Inches(3.2))
    tf = txBox.text_frame
    tf.word_wrap = True
    for j, line in enumerate(lines):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = line
        if line.isupper() and line:  # category header
            run.font.size = Pt(12)
            run.font.color.rgb = GOLD
            run.font.bold = True
        else:
            run.font.size = Pt(11)
            run.font.color.rgb = LIGHT_GRAY
        run.font.name = "Calibri"
        p.line_spacing = Pt(18)


# ── SLIDE 5: Output ─────────────────────────────────────────

s = prs.slides.add_slide(blank_layout)
set_slide_bg(s)
section_label(s, "OUTPUT")
slide_title(s, "What the space can change")

outputs = [
    ("Light", "intensity, color, shadow, emission glow"),
    ("Material", "color, smoothness, opacity, dissolve, rim glow"),
    ("Sound", "volume, pitch, spatial blend, reverb"),
    ("Particles", "emission rate, size, speed, color, burst"),
    ("Environment", "fog, ambient light, skybox"),
    ("Post-Processing", "bloom, saturation, vignette, depth of field"),
    ("Camera", "field of view, shake, position"),
]

for i, (domain, details) in enumerate(outputs):
    y_pos = Inches(2.1 + i * 0.45)
    add_textbox(s, Inches(0.6), y_pos, Inches(2.2), Inches(0.4),
                domain, font_size=13, color=GOLD, bold=True)
    add_textbox(s, Inches(2.8), y_pos, Inches(6.8), Inches(0.4),
                details, font_size=11, color=LIGHT_GRAY)


# ── SLIDE 6: The Transformation ─────────────────────────────

s = prs.slides.add_slide(blank_layout)
set_slide_bg(s)
section_label(s, "TRANSFORMATION")
add_textbox(s, Inches(0.8), Inches(1.2), Inches(8.4), Inches(2.0),
            "Same input.\nSame output.\nCompletely different feel.",
            font_size=36, bold=True)
add_textbox(s, Inches(0.8), Inches(3.8), Inches(8.4), Inches(1.0),
            "The transformation is where design lives.\n"
            "It determines how one becomes the other.",
            font_size=15, color=LIGHT_GRAY, line_spacing=Pt(24))


# ── SLIDE 7: Two Questions ──────────────────────────────────

s = prs.slides.add_slide(blank_layout)
set_slide_bg(s)
section_label(s, "TRANSFORMATION")
slide_title(s, "Two questions define the transformation")

# Question 1
add_rect(s, Inches(0.6), Inches(2.3), Inches(4.2), Inches(2.4),
         RGBColor(0x22, 0x22, 0x3A), BLUE)
add_textbox(s, Inches(0.8), Inches(2.4), Inches(3.8), Inches(0.5),
            "What kind of signal?", font_size=16, bold=True, color=BLUE)
add_textbox(s, Inches(0.8), Inches(3.0), Inches(3.8), Inches(1.5),
            "Continuous\n"
            "a flowing value with magnitude\n"
            "distance, angle, slider position\n"
            "\n"
            "Binary\n"
            "on or off, no in-between\n"
            "contact, button press, toggle",
            font_size=11, color=LIGHT_GRAY, line_spacing=Pt(16))

# Question 2
add_rect(s, Inches(5.2), Inches(2.3), Inches(4.2), Inches(2.4),
         RGBColor(0x22, 0x22, 0x3A), GOLD)
add_textbox(s, Inches(5.4), Inches(2.4), Inches(3.8), Inches(0.5),
            "What kind of relationship?", font_size=16, bold=True, color=GOLD)
add_textbox(s, Inches(5.4), Inches(3.0), Inches(3.8), Inches(1.5),
            "Bound\n"
            "output depends on input\n"
            "when input changes, output follows\n"
            "\n"
            "Unbound\n"
            "output independent after trigger\n"
            "input starts it, then it plays alone",
            font_size=11, color=LIGHT_GRAY, line_spacing=Pt(16))


# ── SLIDE 8: The 4 Modes ────────────────────────────────────

s = prs.slides.add_slide(blank_layout)
set_slide_bg(s)
section_label(s, "THE 4 MODES")
slide_title(s, "Signal \u00d7 Relationship = four interactions", font_size=26)

# Column headers
add_textbox(s, Inches(2.8), Inches(1.9), Inches(3.0), Inches(0.4),
            "BOUND", font_size=14, bold=True, alignment=PP_ALIGN.CENTER)
add_textbox(s, Inches(6.0), Inches(1.9), Inches(3.2), Inches(0.4),
            "UNBOUND", font_size=14, bold=True, alignment=PP_ALIGN.CENTER)

# Row headers
add_textbox(s, Inches(0.3), Inches(2.5), Inches(2.3), Inches(1.2),
            "CONTINUOUS\nsignal", font_size=11, bold=True,
            alignment=PP_ALIGN.RIGHT, line_spacing=Pt(16))
add_textbox(s, Inches(0.3), Inches(3.9), Inches(2.3), Inches(1.2),
            "BINARY\nsignal", font_size=11, bold=True,
            alignment=PP_ALIGN.RIGHT, line_spacing=Pt(16))

modes = [
    # (col, row, name, example, left, top)
    (0, 0, "Tracking",          "Proximity \u2192 Light intensity",
     Inches(2.8), Inches(2.4)),
    (1, 0, "Threshold Trigger", "Gaze duration hits 3s\n\u2192 material dissolves",
     Inches(6.0), Inches(2.4)),
    (0, 1, "Switch",            "Enter zone \u2192\nLight ON / OFF",
     Inches(2.8), Inches(3.8)),
    (1, 1, "One-Shot",          "Button press \u2192\nplay animation",
     Inches(6.0), Inches(3.8)),
]

for col, row, name, example, left, top in modes:
    w = Inches(3.0) if col == 0 else Inches(3.2)
    add_rect(s, left, top, w, Inches(1.2), MODE_BG[name], MODE_COLORS[name], Pt(1.5))
    add_textbox(s, left + Inches(0.15), top + Inches(0.1), w - Inches(0.3), Inches(0.4),
                name, font_size=13, bold=True, color=MODE_COLORS[name])
    add_textbox(s, left + Inches(0.15), top + Inches(0.5), w - Inches(0.3), Inches(0.6),
                example, font_size=10, color=LIGHT_GRAY, line_spacing=Pt(14))


# ── SLIDE 9: Tracking ───────────────────────────────────────

s = prs.slides.add_slide(blank_layout)
set_slide_bg(s)
section_label(s, "TRACKING  \u00b7  Continuous + Bound")
slide_title(s, "Output follows input", font_size=32)
add_textbox(s, Inches(0.6), Inches(2.0), Inches(5.5), Inches(0.5),
            "Proximity \u2192 Light intensity", font_size=18, color=BLUE, bold=True)
slide_body(s,
    "Walk closer, it brightens. Walk away, it dims.\n"
    "The output is always coupled to the input.\n"
    "\n"
    "Design tools:\n"
    "Curve \u2014 the shape of the mapping (linear, eased, inverse)\n"
    "Range \u2014 how much input it listens to, how much output it produces\n"
    "Attack / Release \u2014 how fast it responds, how fast it fades",
    top=Inches(2.7), font_size=13)


# ── SLIDE 10: Switch ────────────────────────────────────────

s = prs.slides.add_slide(blank_layout)
set_slide_bg(s)
section_label(s, "SWITCH  \u00b7  Binary + Bound")
slide_title(s, "Output transitions between two states", font_size=30)
add_textbox(s, Inches(0.6), Inches(2.0), Inches(5.5), Inches(0.5),
            "Enter zone \u2192 Light ON / OFF", font_size=18, color=AMBER, bold=True)
slide_body(s,
    "Inside = ON. Outside = OFF.\n"
    "The output is bound to the binary state.\n"
    "\n"
    "Design tools:\n"
    "Attack \u2014 how fast it turns on (instant snap vs. slow fade in)\n"
    "Release \u2014 how fast it turns off (instant vs. lingering)\n"
    "Asymmetry \u2014 fast on + slow off = the space remembers you",
    top=Inches(2.7), font_size=13)


# ── SLIDE 11: One-Shot ──────────────────────────────────────

s = prs.slides.add_slide(blank_layout)
set_slide_bg(s)
section_label(s, "ONE-SHOT  \u00b7  Binary + Unbound")
slide_title(s, "Input triggers, output plays alone", font_size=30)
add_textbox(s, Inches(0.6), Inches(2.0), Inches(5.5), Inches(0.5),
            "Button press \u2192 play animation", font_size=18, color=RED, bold=True)
slide_body(s,
    "You press the button. The animation plays.\n"
    "It doesn\u2019t matter if you keep pressing or walk away.\n"
    "\n"
    "Design tools:\n"
    "Curve \u2014 the response shape (burst, swell, snap + tail, overshoot)\n"
    "Duration \u2014 how long the response plays\n"
    "The curve IS the temporal behavior \u2014 attack/release are baked in",
    top=Inches(2.7), font_size=13)


# ── SLIDE 12: Threshold Trigger ─────────────────────────────

s = prs.slides.add_slide(blank_layout)
set_slide_bg(s)
section_label(s, "THRESHOLD TRIGGER  \u00b7  Continuous + Unbound")
slide_title(s, "Continuous input builds to a trigger", font_size=30)
add_textbox(s, Inches(0.6), Inches(2.0), Inches(6.5), Inches(0.5),
            "Gaze duration hits 3s \u2192 material dissolves",
            font_size=18, color=GREEN, bold=True)
slide_body(s,
    "Your gaze accumulates over time. You see the number climbing.\n"
    "At 3 seconds, the response fires and plays independently.\n"
    "\n"
    "Design tools:\n"
    "Threshold \u2014 when does the continuous signal trigger? (distance, time, speed)\n"
    "Curve \u2014 the response shape once triggered\n"
    "Duration \u2014 how long the response plays",
    top=Inches(2.7), font_size=13)


# ── SLIDE 13: Shape ─────────────────────────────────────────

s = prs.slides.add_slide(blank_layout)
set_slide_bg(s)
section_label(s, "THE SHAPE OF THE TRANSFORMATION")
slide_title(s, "Three properties define how it feels", font_size=28)

props = [
    ("CURVE", GOLD,
     "The mapping shape. For bound: input value \u2192 output value.\n"
     "For unbound: time \u2192 output value. Same tool, different meaning."),
    ("RANGE", BLUE,
     "Input range: which part of the signal the curve listens to.\n"
     "Output range: how much the output can change. Controls sensitivity."),
    ("ENVELOPE", LIGHT_GRAY,
     "Attack: how fast the output reaches its target.\n"
     "Release: how fast it returns to rest. Together they set the feel."),
]

for i, (name, color, desc) in enumerate(props):
    y = Inches(2.2 + i * 1.0)
    add_textbox(s, Inches(0.6), y, Inches(1.8), Inches(0.4),
                name, font_size=14, bold=True, color=color)
    add_textbox(s, Inches(2.5), y, Inches(7.0), Inches(0.9),
                desc, font_size=11, color=LIGHT_GRAY, line_spacing=Pt(16))


# ── SLIDE 14: Mapping Curves ────────────────────────────────

s = prs.slides.add_slide(blank_layout)
set_slide_bg(s)
section_label(s, "CURVE SHAPES")
slide_title(s, "Mapping curves (Bound)", font_size=26)

add_textbox(s, Inches(0.6), Inches(1.8), Inches(8.8), Inches(0.5),
            "Same input, same output \u2014 the curve changes how one becomes the other.",
            font_size=12, color=LIGHT_GRAY, italic=True)

# Generate and embed curve strip
img_path = make_curve_strip(CURVES, "#6B9EBF", "Mapping", y_max=1.15)
s.shapes.add_picture(img_path, Inches(0.4), Inches(2.5), Inches(9.2))
os.unlink(img_path)


# ── SLIDE 15: Response Curves ───────────────────────────────

s = prs.slides.add_slide(blank_layout)
set_slide_bg(s)
section_label(s, "CURVE SHAPES")
slide_title(s, "Response curves (Unbound)", font_size=26)

add_textbox(s, Inches(0.6), Inches(1.8), Inches(8.8), Inches(0.5),
            "The curve IS the temporal behavior \u2014 it plays out over Duration.",
            font_size=12, color=LIGHT_GRAY, italic=True)

img_path = make_curve_strip(RESPONSE_CURVES, "#E74C3C", "Response", y_max=1.55)
s.shapes.add_picture(img_path, Inches(0.4), Inches(2.5), Inches(9.2))
os.unlink(img_path)


# ── SLIDE 16: Attack & Release ──────────────────────────────

s = prs.slides.add_slide(blank_layout)
set_slide_bg(s)
section_label(s, "ENVELOPE")
slide_title(s, "Attack & Release", font_size=32)

combos = [
    ("Fast attack + Fast release", "Mechanical, crisp", "light switch"),
    ("Fast attack + Slow release", "Eager, lingering", "the space remembers you"),
    ("Slow attack + Fast release", "Reluctant, crisp", "suspicious, then forgets"),
    ("Slow attack + Slow release", "Atmospheric, dreamy", "weather changing"),
]

for i, (combo, feel, metaphor) in enumerate(combos):
    y = Inches(2.1 + i * 0.8)
    add_textbox(s, Inches(0.6), y, Inches(3.5), Inches(0.35),
                combo, font_size=13, bold=True, color=GOLD)
    add_textbox(s, Inches(4.2), y, Inches(2.5), Inches(0.35),
                feel, font_size=13, color=WHITE, italic=True)
    add_textbox(s, Inches(7.0), y, Inches(2.8), Inches(0.35),
                metaphor, font_size=11, color=MID_GRAY)


# ── SLIDE 17: Same atom, different feel ─────────────────────

s = prs.slides.add_slide(blank_layout)
set_slide_bg(s)
section_label(s, "DESIGN IN ACTION")
slide_title(s, "Same input, same output. Different curve.", font_size=26)

add_textbox(s, Inches(0.6), Inches(1.8), Inches(8.8), Inches(0.4),
            "All three: Proximity \u2192 Emission Intensity  (Tracking mode)",
            font_size=12, color=LIGHT_GRAY)

img_path = make_comparison_plot()
s.shapes.add_picture(img_path, Inches(1.5), Inches(2.4), Inches(7.0))
os.unlink(img_path)

add_textbox(s, Inches(0.6), Inches(4.6), Inches(8.8), Inches(0.5),
            "The curve is the personality of the interaction.",
            font_size=14, color=GOLD, italic=True, alignment=PP_ALIGN.CENTER)


# ── SLIDE 18: Coupling Quality ──────────────────────────────

s = prs.slides.add_slide(blank_layout)
set_slide_bg(s)
section_label(s, "COUPLING QUALITY")
slide_title(s, "How does the interaction feel?", font_size=30)

add_textbox(s, Inches(0.6), Inches(1.9), Inches(8.8), Inches(0.4),
            "Use these pairs to describe and design. Not right/wrong \u2014 expressive choices.",
            font_size=11, color=MID_GRAY, italic=True)

qualities = [
    ("Eager", "Reluctant",   "attack speed"),
    ("Crisp", "Lingering",   "release speed"),
    ("Attentive", "Indifferent", "proportionality"),
    ("Precise", "Forgiving",   "range width"),
    ("Alive", "Still",       "oscillation / pulse"),
]

for i, (left_word, right_word, param) in enumerate(qualities):
    y = Inches(2.5 + i * 0.55)
    add_textbox(s, Inches(0.8), y, Inches(2.0), Inches(0.4),
                left_word, font_size=15, bold=True, color=BLUE, alignment=PP_ALIGN.RIGHT)

    # Line
    add_rect(s, Inches(3.0), y + Inches(0.18), Inches(2.5), Inches(0.04),
             MID_GRAY)

    add_textbox(s, Inches(5.7), y, Inches(2.0), Inches(0.4),
                right_word, font_size=15, bold=True, color=RED)

    add_textbox(s, Inches(8.0), y, Inches(1.8), Inches(0.4),
                param, font_size=9, color=DARK_GRAY, italic=True)


# ── SLIDE 19: Close ─────────────────────────────────────────

s = prs.slides.add_slide(blank_layout)
set_slide_bg(s)

add_textbox(s, Inches(0.8), Inches(1.0), Inches(8.4), Inches(1.5),
            "Every interaction is an atom.", font_size=34, bold=True)
add_textbox(s, Inches(0.8), Inches(2.3), Inches(8.4), Inches(1.5),
            "Every atom has a transformation.", font_size=34, bold=True)
add_textbox(s, Inches(0.8), Inches(3.6), Inches(8.4), Inches(1.0),
            "The transformation is what you design.", font_size=34, bold=True, color=GOLD)


# ═══════════════════════════════════════════════════════════════
#  SAVE
# ═══════════════════════════════════════════════════════════════

out_path = (r"D:\Projects Unity\MD4_2026_Mavrodiev_Develop"
            r"\Docs\Interactive Design 4"
            r"\Interaction Model\The Interaction Atom.pptx")
prs.save(out_path)
print(f"Saved: {out_path}")
print(f"Slides: {len(prs.slides)}")
