"""Generate Lecture 4 presentation: Bodies in Space — Interactive Environments"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Colors ──────────────────────────────────────────────────────────────
BG_DARK = RGBColor(0x0D, 0x0D, 0x0D)
BG_SECTION = RGBColor(0x14, 0x14, 0x1E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xBB, 0xBB, 0xBB)
MID_GRAY = RGBColor(0x88, 0x88, 0x88)
ACCENT_WARM = RGBColor(0xFF, 0xA0, 0x40)    # warm orange
ACCENT_COOL = RGBColor(0x40, 0xC0, 0xFF)    # cool cyan
ACCENT_GAZE = RGBColor(0x80, 0xE0, 0xC0)    # mint/teal
ACCENT_CONTACT = RGBColor(0xE0, 0x80, 0xFF)  # purple
ACCENT_PICKUP = RGBColor(0xFF, 0x60, 0x80)   # coral/rose
ACCENT_DIM = RGBColor(0x66, 0x66, 0x66)

prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)

SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height


def set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, left, top, width, height, text, font_size=18,
             color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Segoe UI"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                     Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_para(tf, text, font_size=18, color=WHITE, bold=False,
             space_before=Pt(6), alignment=PP_ALIGN.LEFT, font_name="Segoe UI"):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    if space_before:
        p.space_before = space_before
    return p


def add_rect(slide, left, top, width, height, fill_color, alpha=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


# ════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_bg(s, BG_DARK)

add_text(s, 1.2, 2.0, 13, 1.5, "Bodies in Space",
         font_size=60, color=WHITE, bold=True)
add_text(s, 1.2, 3.6, 13, 1.0, "Interactive Environments",
         font_size=36, color=ACCENT_WARM)
add_text(s, 1.2, 5.2, 10, 0.6, "Lecture 4  \u2014  Interaction Design for Artists",
         font_size=20, color=MID_GRAY)
add_text(s, 1.2, 7.5, 10, 0.5, "NABA Milan  \u00b7  2025\u201326",
         font_size=16, color=ACCENT_DIM)

# ════════════════════════════════════════════════════════════════════════
# SLIDE 2 — The Shift
# ════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, BG_DARK)

add_text(s, 1.2, 1.0, 13, 1.0, "The Shift",
         font_size=44, color=WHITE, bold=True)

# Left column: before
add_rect(s, 1.2, 2.5, 6.2, 4.5, RGBColor(0x1A, 0x1A, 0x1A))
tf = add_text(s, 1.5, 2.7, 5.6, 0.6, "LECTURES 1\u20133", font_size=14, color=MID_GRAY, bold=True)
add_para(tf, "The World Lives On Its Own", font_size=24, color=WHITE, bold=True, space_before=Pt(12))
add_para(tf, "", font_size=8, color=MID_GRAY)
add_para(tf, "Autonomous creatures", font_size=18, color=LIGHT_GRAY)
add_para(tf, "Emergence from simple rules", font_size=18, color=LIGHT_GRAY)
add_para(tf, "Your presence is a disturbance", font_size=18, color=LIGHT_GRAY)
add_para(tf, "", font_size=8, color=MID_GRAY)
add_para(tf, "\u201cWhat happens without me?\u201d", font_size=20, color=MID_GRAY,
         bold=True, alignment=PP_ALIGN.LEFT)

# Right column: after
add_rect(s, 8.6, 2.5, 6.2, 4.5, RGBColor(0x1A, 0x14, 0x0A))
tf = add_text(s, 8.9, 2.7, 5.6, 0.6, "LECTURE 4", font_size=14, color=ACCENT_WARM, bold=True)
add_para(tf, "The World Responds To You", font_size=24, color=WHITE, bold=True, space_before=Pt(12))
add_para(tf, "", font_size=8, color=MID_GRAY)
add_para(tf, "Reactive environments", font_size=18, color=ACCENT_WARM)
add_para(tf, "Body as interface", font_size=18, color=ACCENT_WARM)
add_para(tf, "Your presence is the input", font_size=18, color=ACCENT_WARM)
add_para(tf, "", font_size=8, color=MID_GRAY)
add_para(tf, "\u201cWhat happens because of me?\u201d", font_size=20, color=ACCENT_WARM,
         bold=True, alignment=PP_ALIGN.LEFT)

# ════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Reference: Dune (Roosegaarde)
# ════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, BG_SECTION)

add_text(s, 1.2, 0.8, 10, 0.5, "REFERENCE", font_size=14, color=ACCENT_WARM, bold=True)
add_text(s, 1.2, 1.3, 13, 1.0, "Daan Roosegaarde \u2014 Dune",
         font_size=40, color=WHITE, bold=True)
add_text(s, 1.2, 2.4, 6, 0.4, "2006 \u2013 ongoing  \u00b7  Rotterdam, NL",
         font_size=16, color=MID_GRAY)

tf = add_text(s, 1.2, 3.5, 7, 3.5,
              "Hundreds of light fibers line a public walkway.",
              font_size=22, color=LIGHT_GRAY)
add_para(tf, "", font_size=10, color=MID_GRAY)
add_para(tf, "They glow and hum as you walk past.", font_size=22, color=LIGHT_GRAY)
add_para(tf, "They dim when you leave.", font_size=22, color=LIGHT_GRAY)
add_para(tf, "", font_size=10, color=MID_GRAY)
add_para(tf, "Presence  \u2192  Light + Sound", font_size=26, color=ACCENT_WARM, bold=True)
add_para(tf, "", font_size=10, color=MID_GRAY)
add_para(tf, "One sensor. One response. Hundreds of copies.\nThat\u2019s it. And it feels alive.",
         font_size=18, color=MID_GRAY)

# placeholder for image
add_rect(s, 9.5, 2.5, 5.5, 4.5, RGBColor(0x22, 0x22, 0x22))
add_text(s, 9.5, 4.3, 5.5, 0.5, "[  image: Dune installation  ]",
         font_size=14, color=ACCENT_DIM, alignment=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Reference: Rain Room
# ════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, BG_SECTION)

add_text(s, 1.2, 0.8, 10, 0.5, "REFERENCE", font_size=14, color=ACCENT_COOL, bold=True)
add_text(s, 1.2, 1.3, 13, 1.0, "Random International \u2014 Rain Room",
         font_size=40, color=WHITE, bold=True)
add_text(s, 1.2, 2.4, 6, 0.4, "2012  \u00b7  Barbican, London",
         font_size=16, color=MID_GRAY)

tf = add_text(s, 1.2, 3.5, 7, 3.5,
              "Rain falls everywhere in a dark room.",
              font_size=22, color=LIGHT_GRAY)
add_para(tf, "Except where you stand.", font_size=22, color=WHITE, bold=True)
add_para(tf, "", font_size=10, color=MID_GRAY)
add_para(tf, "Cameras track your position.\nWater nozzles shut off above you.",
         font_size=22, color=LIGHT_GRAY)
add_para(tf, "", font_size=10, color=MID_GRAY)
add_para(tf, "Presence  \u2192  Absence", font_size=26, color=ACCENT_COOL, bold=True)
add_para(tf, "", font_size=10, color=MID_GRAY)
add_para(tf, "The space doesn\u2019t add something.\nIt removes something. That\u2019s still a response.",
         font_size=18, color=MID_GRAY)

add_rect(s, 9.5, 2.5, 5.5, 4.5, RGBColor(0x22, 0x22, 0x22))
add_text(s, 9.5, 4.3, 5.5, 0.5, "[  image: Rain Room  ]",
         font_size=14, color=ACCENT_DIM, alignment=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Reference: Treachery of Sanctuary
# ════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, BG_SECTION)

add_text(s, 1.2, 0.8, 10, 0.5, "REFERENCE", font_size=14, color=ACCENT_CONTACT, bold=True)
add_text(s, 1.2, 1.3, 13, 1.0, "Chris Milk \u2014 The Treachery of Sanctuary",
         font_size=38, color=WHITE, bold=True)
add_text(s, 1.2, 2.4, 6, 0.4, "2012  \u00b7  Creators Project, New York",
         font_size=16, color=MID_GRAY)

tf = add_text(s, 1.2, 3.5, 7, 3.5,
              "Three shadow panels. Same input: your body.",
              font_size=22, color=LIGHT_GRAY)
add_para(tf, "", font_size=10, color=MID_GRAY)
add_para(tf, "Panel 1: your silhouette dissolves into birds", font_size=20, color=LIGHT_GRAY)
add_para(tf, "Panel 2: you grow wings", font_size=20, color=LIGHT_GRAY)
add_para(tf, "Panel 3: you flap your arms, release birds", font_size=20, color=LIGHT_GRAY)
add_para(tf, "", font_size=10, color=MID_GRAY)
add_para(tf, "Body  \u2192  Escalating Transformation", font_size=26, color=ACCENT_CONTACT, bold=True)
add_para(tf, "", font_size=10, color=MID_GRAY)
add_para(tf, "The design isn\u2019t the technology.\nIt\u2019s the mapping between action and response.",
         font_size=18, color=MID_GRAY)

add_rect(s, 9.5, 2.5, 5.5, 4.5, RGBColor(0x22, 0x22, 0x22))
add_text(s, 9.5, 4.3, 5.5, 0.5, "[  image: Treachery of Sanctuary  ]",
         font_size=14, color=ACCENT_DIM, alignment=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════════════
# SLIDE 6 — The Question
# ════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, BG_DARK)

add_text(s, 1.2, 2.5, 13.5, 2.0,
         "What does the space\nknow about you?",
         font_size=54, color=WHITE, bold=True)
add_text(s, 1.2, 5.5, 13.5, 1.0,
         "And what does it do with that knowledge?",
         font_size=30, color=ACCENT_WARM)

# ════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Four Inputs
# ════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, BG_DARK)

add_text(s, 1.2, 0.6, 13, 0.8, "How the Body Speaks",
         font_size=40, color=WHITE, bold=True)
add_text(s, 1.2, 1.3, 13, 0.5, "Four things the space can sense about you",
         font_size=18, color=MID_GRAY)

# Four cards
card_data = [
    ("PRESENCE", "You\u2019re near something", "ProximitySensor", ACCENT_WARM, 1.2),
    ("GAZE", "You\u2019re looking at something", "GazeSensor", ACCENT_GAZE, 4.8),
    ("CONTACT", "You\u2019re touching something", "TriggerSensor", ACCENT_CONTACT, 8.4),
    ("INTERACTION", "You picked something up", "Pickup (new)", ACCENT_PICKUP, 12.0),
]

for title, desc, tech, color, left in card_data:
    add_rect(s, left, 2.3, 3.0, 5.0, RGBColor(0x1A, 0x1A, 0x1A))
    # colored top accent bar
    add_rect(s, left, 2.3, 3.0, 0.08, color)
    add_text(s, left + 0.3, 2.7, 2.4, 0.5, title,
             font_size=18, color=color, bold=True)
    add_text(s, left + 0.3, 3.4, 2.4, 1.5, desc,
             font_size=20, color=WHITE)
    add_text(s, left + 0.3, 5.8, 2.4, 0.5, tech,
             font_size=13, color=MID_GRAY, font_name="Consolas")
    # icon placeholder
    add_rect(s, left + 0.8, 4.5, 1.4, 1.0, RGBColor(0x22, 0x22, 0x22))

# ════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Four Outputs
# ════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, BG_DARK)

add_text(s, 1.2, 0.6, 13, 0.8, "How the Space Answers",
         font_size=40, color=WHITE, bold=True)
add_text(s, 1.2, 1.3, 13, 0.5, "Four ways the environment can respond",
         font_size=18, color=MID_GRAY)

out_data = [
    ("LIGHT", "Brightness, color,\nshadow, glow", "Light + Emission", ACCENT_WARM, 1.2),
    ("SURFACE", "Material color,\ntexture, transparency", "MaterialColor\nMaterialFloat", ACCENT_GAZE, 4.8),
    ("SOUND", "Spatial audio,\nambient tones", "SoundPlayer (new)", ACCENT_CONTACT, 8.4),
    ("ATMOSPHERE", "Fog, post-processing,\nskybox, mood", "AtmosphereControl\n(new)", ACCENT_PICKUP, 12.0),
]

for title, desc, tech, color, left in out_data:
    add_rect(s, left, 2.3, 3.0, 5.0, RGBColor(0x1A, 0x1A, 0x1A))
    add_rect(s, left, 2.3, 3.0, 0.08, color)
    add_text(s, left + 0.3, 2.7, 2.4, 0.5, title,
             font_size=18, color=color, bold=True)
    add_text(s, left + 0.3, 3.4, 2.4, 1.5, desc,
             font_size=20, color=WHITE)
    add_text(s, left + 0.3, 5.8, 2.4, 0.5, tech,
             font_size=13, color=MID_GRAY, font_name="Consolas")
    add_rect(s, left + 0.8, 4.5, 1.4, 1.0, RGBColor(0x22, 0x22, 0x22))

# ════════════════════════════════════════════════════════════════════════
# SLIDE 9 — The Matrix
# ════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, BG_DARK)

add_text(s, 1.2, 0.6, 13, 0.8, "The Design Space",
         font_size=40, color=WHITE, bold=True)
add_text(s, 1.2, 1.3, 13, 0.5, "Any input can drive any output \u2014 the mapping is the design",
         font_size=18, color=MID_GRAY)

# Table grid
grid_left = 2.5
grid_top = 2.5
col_w = 2.5
row_h = 1.1
headers_col = ["Light", "Surface", "Sound", "Atmosphere"]
headers_row = ["Presence", "Gaze", "Contact", "Pickup"]
header_colors = [ACCENT_WARM, ACCENT_GAZE, ACCENT_CONTACT, ACCENT_PICKUP]

# Column headers
for i, h in enumerate(headers_col):
    add_text(s, grid_left + 2.5 + i * col_w, grid_top - 0.5, col_w, 0.5, h,
             font_size=16, color=ACCENT_WARM if i == 0 else
             ACCENT_GAZE if i == 1 else ACCENT_CONTACT if i == 2 else ACCENT_PICKUP,
             bold=True, alignment=PP_ALIGN.CENTER)

# Row headers + cells
cell_content = [
    ["Pillar glow", "?", "?", "?"],
    ["?", "Panel reveal", "?", "?"],
    ["?", "Tile color", "Tile sound", "?"],
    ["?", "?", "?", "Room transform"],
]

for r, (row_label, row_color) in enumerate(zip(headers_row, header_colors)):
    y = grid_top + r * row_h
    add_text(s, grid_left - 0.3, y, 2.5, row_h, row_label,
             font_size=18, color=row_color, bold=True, alignment=PP_ALIGN.RIGHT)

    for c in range(4):
        x = grid_left + 2.5 + c * col_w
        content = cell_content[r][c]
        is_filled = content != "?"
        cell_bg = RGBColor(0x2A, 0x20, 0x10) if is_filled else RGBColor(0x18, 0x18, 0x18)
        cell_text_color = WHITE if is_filled else RGBColor(0x44, 0x44, 0x44)
        add_rect(s, x + 0.05, y + 0.05, col_w - 0.1, row_h - 0.1, cell_bg)
        add_text(s, x + 0.1, y + 0.15, col_w - 0.2, row_h - 0.3, content,
                 font_size=15, color=cell_text_color, alignment=PP_ALIGN.CENTER)

add_text(s, 2.5, 7.2, 10, 0.6,
         "16 possible mappings. Today we build 4. You choose the rest.",
         font_size=20, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════════════
# SLIDE 10 — The Spec-Driven Approach
# ════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, BG_DARK)

add_text(s, 1.2, 0.6, 13, 0.8, "From Wiring to Specifying",
         font_size=40, color=WHITE, bold=True)
add_text(s, 1.2, 1.3, 13, 0.5, "A new way to create behaviors",
         font_size=18, color=MID_GRAY)

# Left: before
add_rect(s, 1.2, 2.3, 6.5, 5.5, RGBColor(0x1A, 0x1A, 0x1A))
tf = add_text(s, 1.5, 2.5, 6, 0.5, "BEFORE: Inspector Wiring",
              font_size=16, color=MID_GRAY, bold=True)
add_para(tf, "", font_size=8, color=MID_GRAY)
add_para(tf, "[Pillar]", font_size=16, color=ACCENT_DIM, font_name="Consolas")
add_para(tf, "  \u251c\u2500 ProximitySensor", font_size=16, color=LIGHT_GRAY, font_name="Consolas")
add_para(tf, "  \u251c\u2500 SensorResponse", font_size=16, color=LIGHT_GRAY, font_name="Consolas")
add_para(tf, "  \u251c\u2500 MaterialEmissionIntensity (on)", font_size=16, color=LIGHT_GRAY, font_name="Consolas")
add_para(tf, "  \u251c\u2500 MaterialEmissionIntensity (off)", font_size=16, color=LIGHT_GRAY, font_name="Consolas")
add_para(tf, "  \u2514\u2500 ??? (no light control)", font_size=16, color=RGBColor(0xFF, 0x55, 0x55), font_name="Consolas")
add_para(tf, "", font_size=10, color=MID_GRAY)
add_para(tf, "5 components, 6 event wires, can\u2019t control light", font_size=16, color=MID_GRAY)

# Right: after
add_rect(s, 8.3, 2.3, 6.5, 5.5, RGBColor(0x1A, 0x14, 0x0A))
tf = add_text(s, 8.6, 2.5, 6, 0.5, "AFTER: Spec-Driven Controller",
              font_size=16, color=ACCENT_WARM, bold=True)
add_para(tf, "", font_size=8, color=MID_GRAY)
add_para(tf, "[Pillar]", font_size=16, color=ACCENT_DIM, font_name="Consolas")
add_para(tf, "  \u251c\u2500 ProximitySensor", font_size=16, color=ACCENT_WARM, font_name="Consolas")
add_para(tf, "  \u2514\u2500 GlowOnPresence", font_size=16, color=ACCENT_WARM, font_name="Consolas", bold=True)
add_para(tf, "", font_size=10, color=MID_GRAY)
add_para(tf, "2 components. Name = behavior.", font_size=16, color=ACCENT_WARM)
add_para(tf, "Handles emission + light + animation.", font_size=16, color=ACCENT_WARM)
add_para(tf, "", font_size=14, color=MID_GRAY)
add_para(tf, "You describe what you want.", font_size=20, color=WHITE, bold=True)
add_para(tf, "AI writes the code.", font_size=20, color=WHITE, bold=True)
add_para(tf, "You wire 2 references.", font_size=20, color=WHITE, bold=True)

# ════════════════════════════════════════════════════════════════════════
# SLIDE 11 — What a Spec Looks Like
# ════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, BG_DARK)

add_text(s, 1.2, 0.6, 13, 0.8, "The Spec",
         font_size=40, color=WHITE, bold=True)
add_text(s, 1.2, 1.3, 13, 0.5, "Plain language \u2192 describe references, config, behavior \u2192 AI generates the script",
         font_size=18, color=MID_GRAY)

add_rect(s, 1.2, 2.2, 13.5, 5.8, RGBColor(0x14, 0x14, 0x14))

spec_lines = [
    ("# GlowOnPresence", ACCENT_WARM, True, 16),
    ("", MID_GRAY, False, 8),
    ("## References", ACCENT_GAZE, True, 15),
    ("- ProximitySensor (on same object)", LIGHT_GRAY, False, 15),
    ("- Renderer (mesh renderer)", LIGHT_GRAY, False, 15),
    ("- Light (child point light)", LIGHT_GRAY, False, 15),
    ("", MID_GRAY, False, 8),
    ("## Config", ACCENT_GAZE, True, 15),
    ("- emissionColor (Color, default warm orange)", LIGHT_GRAY, False, 15),
    ("- activeIntensity (float, default 3.0)", LIGHT_GRAY, False, 15),
    ("- idleIntensity (float, default 0.0)", LIGHT_GRAY, False, 15),
    ("- fadeSpeed (float, default 0.5)", LIGHT_GRAY, False, 15),
    ("", MID_GRAY, False, 8),
    ("## Behavior", ACCENT_GAZE, True, 15),
    ("- When sensor detects: animate emission + light to active values", LIGHT_GRAY, False, 15),
    ("- When sensor loses all: animate back to idle values", LIGHT_GRAY, False, 15),
    ("- Use DOTween for smooth animation", LIGHT_GRAY, False, 15),
]

tf = add_text(s, 1.6, 2.4, 12.5, 0.3, spec_lines[0][0],
              font_size=spec_lines[0][3], color=spec_lines[0][1],
              bold=spec_lines[0][2], font_name="Consolas")
for text, color, bold, size in spec_lines[1:]:
    add_para(tf, text, font_size=size, color=color, bold=bold,
             space_before=Pt(2), font_name="Consolas")

# ════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Today's Build
# ════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, BG_DARK)

add_text(s, 1.2, 0.6, 13, 0.8, "Today\u2019s Build: The Responsive Room",
         font_size=40, color=WHITE, bold=True)

steps = [
    ("1", "Presence \u2192 Light", "Pillars glow when you walk near",
     "ProximitySensor + GlowOnPresence", ACCENT_WARM),
    ("2", "Gaze \u2192 Surface", "Wall panels reveal when you look at them",
     "GazeSensor + RevealOnGaze", ACCENT_GAZE),
    ("3", "Contact \u2192 Sound", "Floor tiles play tones when you step on them",
     "TriggerSensor + SoundOnContact", ACCENT_CONTACT),
    ("4", "Pickup \u2192 Atmosphere", "Crystals transform the entire room",
     "Pickup + AtmosphereControl", ACCENT_PICKUP),
]

for i, (num, title, desc, tech, color) in enumerate(steps):
    y = 1.8 + i * 1.7
    # number circle
    add_rect(s, 1.2, y, 0.7, 0.7, color)
    add_text(s, 1.2, y + 0.05, 0.7, 0.7, num,
             font_size=28, color=BG_DARK, bold=True, alignment=PP_ALIGN.CENTER)
    # text
    add_text(s, 2.3, y, 5, 0.5, title,
             font_size=24, color=color, bold=True)
    add_text(s, 2.3, y + 0.5, 6, 0.5, desc,
             font_size=18, color=LIGHT_GRAY)
    add_text(s, 9.5, y + 0.15, 5, 0.5, tech,
             font_size=14, color=MID_GRAY, font_name="Consolas")

add_text(s, 1.2, 8.0, 13, 0.5,
         "One room. Four layers. Each layer is a spec you write \u2192 AI generates \u2192 you wire.",
         font_size=18, color=MID_GRAY)

# ════════════════════════════════════════════════════════════════════════
# SLIDE 13 (bonus) — The Takeaway
# ════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, BG_DARK)

add_text(s, 1.2, 1.5, 13.5, 2.0,
         "The mapping is the art.",
         font_size=54, color=WHITE, bold=True)

tf = add_text(s, 1.2, 4.0, 13.5, 1.0,
              "Same room. Same walls. Same pillars.",
              font_size=28, color=MID_GRAY)
add_para(tf, "But what the space knows about you \u2014 and what it does with that \u2014",
         font_size=28, color=MID_GRAY)
add_para(tf, "changes everything.", font_size=28, color=ACCENT_WARM, bold=True)

add_text(s, 1.2, 7.0, 13.5, 1.0,
         "Dune fills one cell in the matrix and makes it extraordinary.\nYou don\u2019t need sixteen. You need the right one.",
         font_size=20, color=ACCENT_DIM)

# ════════════════════════════════════════════════════════════════════════
# Save
# ════════════════════════════════════════════════════════════════════════
output_path = r"D:\Projects Unity\Ludocore\Documentation\Courses\Interactive Design 4\Lecture 4\Lecture 4 - Bodies in Space.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")
print(f"Slides: {len(prs.slides)}")
